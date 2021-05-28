from pptx import Presentation
from pptx.util import Cm
from pathlib import Path
import argparse
  
if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='Deduplicate consecutive images')
    parser.add_argument('--width_pixel', required=True, default=480, type=int,
            help='Set the width of the video in pixels')
    parser.add_argument('--height_pixel', required=True, default=270, type=int,
            help='Set the height of the video in pixels')
    parser.add_argument('--pixels_per_cm', required=False, default=56.69291338582677,
            help='Number of pixels per centimeter')
    args = parser.parse_args()

    # Creating presentation object
    presentation = Presentation()
    width_cm = (1.0/args.pixels_per_cm) * args.width_pixel
    height_cm = (1.0/args.pixels_per_cm) * args.height_pixel
    presentation.slide_width = Cm(width_cm)
    presentation.slide_height = Cm(height_cm)

    print(f'Slide width: {width_cm} cm')
    print(f'Slide height: {height_cm} cm')
      
    # Creating slide layout
    first_slide_layout = presentation.slide_layouts[6] 
      
    """ Ref for slide types: 
    0 ->  title and subtitle
    1 ->  title and content
    2 ->  section header
    3 ->  two content
    4 ->  Comparison
    5 ->  Title only 
    6 ->  Blank
    7 ->  Content with caption
    8 ->  Pic with caption
    """

    left = top = Cm(0)
    width = presentation.slide_width 
    height = presentation.slide_height
      
    for img_path in sorted(Path(r'./').glob('frame_*.jpg'), key=lambda path: float(path.stem.rsplit("_", 1)[1])):
        slide = presentation.slides.add_slide(first_slide_layout)
        pic = slide.shapes.add_picture(str(img_path), left, top, width=width, height=height)
      
    presentation.save("story.pptx")
